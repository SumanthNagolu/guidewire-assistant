#!/usr/bin/env python3
"""
GUIDEWIRE GURU - CONTENT EXTRACTION SCRIPT
==========================================
Extracts text content from PDF, PPTX, DOCX, TXT, MD, and code files.
Saves extracted content as JSON files for embedding pipeline.

Usage:
    python scripts/extract-content.py <input_dir> <output_dir>

Example:
    python scripts/extract-content.py ./guidewire-knowledge ./extracted-knowledge

Requirements:
    pip install python-pptx PyPDF2 python-docx
"""

import os
import json
import sys
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    import PyPDF2
    from docx import Document
except ImportError:
    print("⚠️  Missing dependencies. Installing...")
    print("Run: pip install python-pptx PyPDF2 python-docx")
    sys.exit(1)


def extract_pptx(file_path):
    """Extract text from PowerPoint files"""
    try:
        prs = Presentation(file_path)
        slides_text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = f"=== Slide {slide_num} ===\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content += shape.text + "\n"
            
            if slide_content.strip() != f"=== Slide {slide_num} ===":
                slides_text.append(slide_content)
        
        return "\n\n".join(slides_text)
    except Exception as e:
        raise Exception(f"PPT extraction error: {str(e)}")


def extract_pdf(file_path):
    """Extract text from PDF files"""
    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            pages_text = []
            
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    pages_text.append(f"=== Page {page_num} ===\n{text}")
            
            return "\n\n".join(pages_text)
    except Exception as e:
        raise Exception(f"PDF extraction error: {str(e)}")


def extract_docx(file_path):
    """Extract text from Word documents"""
    try:
        doc = Document(file_path)
        paragraphs = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        
        return "\n\n".join(paragraphs)
    except Exception as e:
        raise Exception(f"DOCX extraction error: {str(e)}")


def extract_text(file_path):
    """Extract text from plain text files"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        raise Exception(f"Text extraction error: {str(e)}")


def detect_source_type(file_path):
    """Detect source type from file path"""
    path_lower = file_path.lower()
    
    if 'resume' in path_lower or 'cv' in path_lower:
        return 'resume'
    elif 'code' in path_lower or 'example' in path_lower or any(path_lower.endswith(ext) for ext in ['.java', '.js', '.py', '.ts']):
        return 'code'
    elif 'interview' in path_lower or 'question' in path_lower:
        return 'interview'
    elif 'project' in path_lower or 'assignment' in path_lower:
        return 'project'
    else:
        return 'guidewire_doc'


def detect_product(content):
    """Detect Guidewire product from content"""
    content_lower = content.lower()
    
    # Check for specific products
    if 'claimcenter' in content_lower or 'claim center' in content_lower:
        return 'ClaimCenter'
    elif 'policycenter' in content_lower or 'policy center' in content_lower:
        return 'PolicyCenter'
    elif 'billingcenter' in content_lower or 'billing center' in content_lower:
        return 'BillingCenter'
    elif 'producerengage' in content_lower or 'producer engage' in content_lower:
        return 'ProducerEngage'
    elif 'customerengage' in content_lower or 'customer engage' in content_lower:
        return 'CustomerEngage'
    elif 'guidewire' in content_lower:
        return 'General'
    else:
        return 'General'


def detect_difficulty(content):
    """Detect difficulty level from content"""
    content_lower = content.lower()
    
    if 'beginner' in content_lower or 'introduction' in content_lower or 'basic' in content_lower:
        return 'beginner'
    elif 'advanced' in content_lower or 'expert' in content_lower or 'senior' in content_lower:
        return 'advanced'
    elif 'intermediate' in content_lower:
        return 'intermediate'
    else:
        return None


def clean_text(text):
    """Clean extracted text"""
    # Remove excessive whitespace
    text = '\n'.join([line.strip() for line in text.split('\n')])
    # Remove more than 2 consecutive newlines
    while '\n\n\n' in text:
        text = text.replace('\n\n\n', '\n\n')
    return text.strip()


def process_file(file_path, output_dir):
    """Process a single file and save as JSON"""
    file_name = os.path.basename(file_path)
    file_ext = os.path.splitext(file_name)[1].lower()
    
    # Extract content based on file type
    try:
        if file_ext == '.pptx':
            content = extract_pptx(file_path)
        elif file_ext == '.pdf':
            content = extract_pdf(file_path)
        elif file_ext == '.docx':
            content = extract_docx(file_path)
        elif file_ext in ['.txt', '.md', '.java', '.js', '.py', '.ts', '.jsx', '.tsx', '.gosu']:
            content = extract_text(file_path)
        else:
            return None  # Skip unsupported formats
        
        if not content or len(content.strip()) < 50:
            return None  # Skip empty or very short files
        
        # Clean content
        content = clean_text(content)
        
        # Detect metadata
        source_type = detect_source_type(file_path)
        product = detect_product(content)
        difficulty = detect_difficulty(content)
        
        # Create output object
        output = {
            'file_name': file_name,
            'file_path': file_path,
            'source_type': source_type,
            'product': product,
            'difficulty': difficulty,
            'content': content,
            'word_count': len(content.split()),
            'extracted_at': datetime.now().isoformat()
        }
        
        # Save as JSON
        output_file = os.path.join(output_dir, f"{Path(file_name).stem}.json")
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(output, f, indent=2, ensure_ascii=False)
        
        return output_file
        
    except Exception as e:
        print(f"  ✗ Error processing {file_name}: {str(e)}")
        return None


def process_directory(input_dir, output_dir):
    """Process all files in directory recursively"""
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    
    print(f"\n🚀 Starting extraction from: {input_dir}")
    print(f"📁 Output directory: {output_dir}\n")
    
    stats = {
        'total': 0,
        'success': 0,
        'failed': 0,
        'skipped': 0,
        'by_type': {}
    }
    
    # Walk through all files
    for root, dirs, files in os.walk(input_dir):
        for file in files:
            file_path = os.path.join(root, file)
            file_ext = os.path.splitext(file)[1].lower()
            
            # Track by type
            if file_ext not in stats['by_type']:
                stats['by_type'][file_ext] = 0
            
            stats['total'] += 1
            
            # Skip hidden files
            if file.startswith('.'):
                stats['skipped'] += 1
                continue
            
            # Process file
            print(f"Processing: {file}...", end=' ')
            result = process_file(file_path, output_dir)
            
            if result:
                stats['success'] += 1
                stats['by_type'][file_ext] += 1
                print("✓")
            else:
                stats['failed'] += 1
                print("✗ (skipped)")
    
    # Print summary
    print("\n" + "="*60)
    print("📊 EXTRACTION SUMMARY")
    print("="*60)
    print(f"Total files found:     {stats['total']}")
    print(f"Successfully extracted: {stats['success']}")
    print(f"Failed/Skipped:        {stats['failed'] + stats['skipped']}")
    print(f"\nBy file type:")
    for ext, count in sorted(stats['by_type'].items()):
        if count > 0:
            print(f"  {ext:10s}: {count}")
    print("="*60)
    print(f"\n✅ Extracted files saved to: {output_dir}")


def main():
    if len(sys.argv) < 3:
        print("Usage: python extract-content.py <input_dir> <output_dir>")
        print("Example: python extract-content.py ./guidewire-knowledge ./extracted-knowledge")
        sys.exit(1)
    
    input_dir = sys.argv[1]
    output_dir = sys.argv[2]
    
    if not os.path.exists(input_dir):
        print(f"❌ Error: Input directory '{input_dir}' does not exist")
        sys.exit(1)
    
    process_directory(input_dir, output_dir)


if __name__ == "__main__":
    main()
