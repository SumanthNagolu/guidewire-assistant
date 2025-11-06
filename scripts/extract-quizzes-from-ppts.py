#!/usr/bin/env python3
"""
Bulk extract quiz questions from PowerPoint files.

This script:
1. Scans all PPT files in data/ folder
2. Finds "Lesson objectives review" slide
3. Extracts all slides after it (quiz questions/answers)
4. Parses question/answer pairs
5. Generates SQL INSERT statements

Usage:
    python3 scripts/extract-quizzes-from-ppts.py

Requirements:
    pip install python-pptx

Output:
    database/BULK-QUIZ-INSERTS.sql
"""

import re
import json
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("❌ Error: python-pptx not installed")
    print("   Run: pip install python-pptx")
    exit(1)


def extract_slide_text(slide) -> str:
    """Extract all text from a slide."""
    text_parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text_parts.append(shape.text)
    return "\n".join(text_parts)


def get_slide_background_color(slide) -> tuple:
    """Try to determine slide background color (RGB)."""
    try:
        if slide.background.fill.type == 1:  # Solid fill
            color = slide.background.fill.fore_color
            if hasattr(color, 'rgb'):
                return (color.rgb[0], color.rgb[1], color.rgb[2])
    except:
        pass
    return None


def is_lesson_review_slide(slide_text: str) -> bool:
    """Check if slide is 'Lesson objectives review'."""
    normalized = slide_text.lower().strip()
    return (
        "lesson objectives review" in normalized or
        "lesson objective review" in normalized or
        "objectives review" in normalized
    )


def is_question_slide(slide_text: str) -> bool:
    """Heuristic: slide with 'Question' header or multiple choice options."""
    text = slide_text.strip()
    # Look for "Question" at start or numbered questions
    if re.match(r'^question\s*\d*', text, re.IGNORECASE):
        return True
    # Look for A) B) C) D) pattern
    if re.search(r'[A-D]\)', text):
        return True
    return False


def is_answer_slide(slide_text: str) -> bool:
    """Heuristic: slide with 'Answer' header."""
    text = slide_text.lower().strip()
    return text.startswith('answer') or 'correct answer' in text


def parse_question_slide(slide_text: str) -> dict:
    """Parse a question slide into structured data."""
    lines = [line.strip() for line in slide_text.split('\n') if line.strip()]
    
    # First non-empty line after "Question" is the question text
    question_text = ''
    options = {}
    
    question_started = False
    for line in lines:
        if re.match(r'^question\s*\d*', line, re.IGNORECASE):
            question_started = True
            continue
        
        if not question_started:
            continue
        
        # Check for options A) B) C) D)
        option_match = re.match(r'([A-D])\)\s*(.+)', line, re.IGNORECASE)
        if option_match:
            letter, text = option_match.groups()
            options[letter.upper()] = text.strip()
        elif not options and question_text:
            # Multi-line question
            question_text += ' ' + line
        elif not options:
            # First line of question
            question_text = line
    
    return {
        'question': question_text.strip(),
        'options': options
    }


def parse_answer_slide(slide_text: str) -> str:
    """Parse an answer slide to get the correct answer letter."""
    text = slide_text.strip()
    
    # Look for patterns like "Answer: B" or "B)" or just "B"
    match = re.search(r'answer[:\s]*([A-D])\)?', text, re.IGNORECASE)
    if match:
        return match.group(1).upper()
    
    # Look for standalone letter
    match = re.search(r'^([A-D])\)?$', text, re.MULTILINE | re.IGNORECASE)
    if match:
        return match.group(1).upper()
    
    # Look for "A supported web browser..." pattern (full answer text)
    for letter in ['A', 'B', 'C', 'D']:
        if f'{letter})' in text or f'{letter} )' in text:
            return letter
    
    return None


def extract_quiz_from_ppt(ppt_path: Path) -> dict:
    """Extract quiz questions from a single PPT file."""
    print(f"  📄 Processing: {ppt_path.name}")
    
    try:
        prs = Presentation(str(ppt_path))
    except Exception as e:
        print(f"     ❌ Error opening PPT: {e}")
        return None
    
    slides = list(prs.slides)
    
    # Find "Lesson objectives review" slide
    review_slide_idx = None
    for idx, slide in enumerate(slides):
        text = extract_slide_text(slide)
        if is_lesson_review_slide(text):
            review_slide_idx = idx
            print(f"     ✅ Found review slide at position {idx + 1}")
            break
    
    if review_slide_idx is None:
        print(f"     ⚠️  No 'Lesson objectives review' found")
        return None
    
    # Extract slides after review
    quiz_slides = slides[review_slide_idx + 1:]
    
    if not quiz_slides:
        print(f"     ⚠️  No slides after review")
        return None
    
    # Parse question/answer pairs
    questions = []
    i = 0
    while i < len(quiz_slides):
        slide = quiz_slides[i]
        slide_text = extract_slide_text(slide)
        
        if is_question_slide(slide_text):
            question_data = parse_question_slide(slide_text)
            
            # Look for answer in next slide
            correct_answer = None
            if i + 1 < len(quiz_slides):
                answer_text = extract_slide_text(quiz_slides[i + 1])
                if is_answer_slide(answer_text):
                    correct_answer = parse_answer_slide(answer_text)
                    i += 1  # Skip answer slide
            
            if question_data['question'] and question_data['options'] and correct_answer:
                questions.append({
                    'question': question_data['question'],
                    'options': question_data['options'],
                    'correct_answer': correct_answer
                })
        
        i += 1
    
    if questions:
        print(f"     ✅ Extracted {len(questions)} questions")
        return {
            'ppt_name': ppt_path.stem,
            'questions': questions
        }
    else:
        print(f"     ⚠️  No questions extracted")
        return None


def map_ppt_to_topic_code(ppt_name: str) -> str:
    """Map PPT filename to topic code."""
    # Try to match patterns like "IS_Claim_01" → "cc-01-001"
    
    # ClaimCenter patterns
    if 'claim' in ppt_name.lower() or 'is_claim' in ppt_name.lower():
        match = re.search(r'(\d{2})', ppt_name)
        if match:
            num = match.group(1)
            return f"cc-01-{num}1"  # cc-01-011, cc-01-021, etc.
    
    # PolicyCenter patterns
    if 'policy' in ppt_name.lower() or 'pp_' in ppt_name.lower():
        match = re.search(r'(\d{2})', ppt_name)
        if match:
            num = match.group(1)
            return f"pc-02-{num}1"
    
    # BillingCenter patterns
    if 'billing' in ppt_name.lower() or 'bc_' in ppt_name.lower():
        match = re.search(r'(\d{2})', ppt_name)
        if match:
            num = match.group(1)
            return f"bc-01-{num}1"
    
    # Foundation/Common patterns
    if 'chapter' in ppt_name.lower() or 'common' in ppt_name.lower():
        match = re.search(r'(\d{1,2})', ppt_name)
        if match:
            num = match.group(1).zfill(2)
            return f"fw-01-{num}1"
    
    # Default: generate from filename
    return f"unknown-{ppt_name[:10]}"


def generate_sql_for_quiz(topic_code: str, quiz_data: dict) -> list:
    """Generate SQL statements for a single quiz."""
    
    def escape_sql(text: str) -> str:
        if text is None:
            return ''
        return text.replace("'", "''")
    
    sql_lines = [
        f"-- Quiz for: {quiz_data['ppt_name']}",
        f"-- Topic code: {topic_code}",
        f"-- Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
        "",
        "DO $$",
        "DECLARE",
        "  v_quiz_id UUID;",
        "BEGIN",
        "  -- Create quiz",
        "  INSERT INTO quizzes (",
        "    id, topic_id, title, description, passing_score, time_limit_minutes, published",
        "  ) VALUES (",
        "    gen_random_uuid(),",
        f"    (SELECT id FROM topics WHERE code = '{topic_code}'),",
        f"    '{escape_sql(quiz_data['ppt_name'])} - Knowledge Check',",
        f"    'Quiz extracted from {escape_sql(quiz_data['ppt_name'])}',",
        "    70,",
        "    15,",
        "    true",
        "  ) RETURNING id INTO v_quiz_id;",
        "",
        "  -- Create questions",
    ]
    
    for idx, q in enumerate(quiz_data['questions'], 1):
        question_text = escape_sql(q['question'])
        options_json = json.dumps(q['options']).replace("'", "''")
        
        sql_lines.extend([
            f"  -- Question {idx}",
            "  INSERT INTO quiz_questions (",
            "    id, quiz_id, position, question, options, correct_answer, points",
            "  ) VALUES (",
            "    gen_random_uuid(),",
            "    v_quiz_id,",
            f"    {idx},",
            f"    '{question_text}',",
            f"    '{options_json}'::jsonb,",
            f"    '{q['correct_answer']}',",
            "    1",
            "  );",
            "",
        ])
    
    sql_lines.extend([
        f"  RAISE NOTICE 'Created quiz for {topic_code} with % questions', {len(quiz_data['questions'])};",
        "END $$;",
        "",
        "",
    ])
    
    return sql_lines


def main():
    print("🚀 Bulk Quiz Extraction from PPT Files")
    print("=" * 60)
    
    # Find all PPT files
    data_dir = Path("data")
    ppt_files = list(data_dir.rglob("*.pptx")) + list(data_dir.rglob("*.ppt"))
    
    if not ppt_files:
        print("❌ No PPT files found in data/ folder")
        return
    
    print(f"\n📁 Found {len(ppt_files)} PPT files\n")
    
    # Extract quizzes
    all_quizzes = []
    for ppt_file in ppt_files:
        quiz_data = extract_quiz_from_ppt(ppt_file)
        if quiz_data:
            topic_code = map_ppt_to_topic_code(ppt_file.stem)
            all_quizzes.append((topic_code, quiz_data))
    
    if not all_quizzes:
        print("\n❌ No quizzes extracted from any PPT file")
        print("   Check if PPTs have 'Lesson objectives review' slide")
        return
    
    print(f"\n✅ Successfully extracted {len(all_quizzes)} quizzes")
    print("\n📝 Generating SQL...")
    
    # Generate SQL
    all_sql = [
        "-- Bulk Quiz Import",
        f"-- Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
        f"-- Total quizzes: {len(all_quizzes)}",
        "",
        "-- Note: This will create quizzes for all topics that have matching codes",
        "-- Topics without matching codes will be skipped",
        "",
        "",
    ]
    
    for topic_code, quiz_data in all_quizzes:
        all_sql.extend(generate_sql_for_quiz(topic_code, quiz_data))
    
    # Add verification query
    all_sql.extend([
        "-- Verification: List all quizzes created",
        "SELECT ",
        "  t.code,",
        "  t.title,",
        "  q.title as quiz_title,",
        "  COUNT(qq.id) as question_count",
        "FROM quizzes q",
        "JOIN topics t ON q.topic_id = t.id",
        "LEFT JOIN quiz_questions qq ON qq.quiz_id = q.id",
        "GROUP BY t.code, t.title, q.title",
        "ORDER BY t.code;",
    ])
    
    # Write to file
    output_file = Path("database/BULK-QUIZ-INSERTS.sql")
    output_file.write_text('\n'.join(all_sql))
    
    print(f"✅ Generated: {output_file}")
    print("\n📋 Next steps:")
    print("   1. Review the generated SQL file")
    print("   2. Open Supabase SQL Editor")
    print(f"   3. Copy and paste contents of {output_file}")
    print("   4. Run the script")
    print("   5. Check verification query at the end")
    print(f"\n🎉 {len(all_quizzes)} quizzes ready to import!")


if __name__ == '__main__':
    main()

